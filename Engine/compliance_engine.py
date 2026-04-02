import re
import time
import os
import difflib
from datetime import datetime, timezone
from typing import Dict, Any, List

# Document parsers
import PyPDF2
import docx
import openpyxl
import pptx

class ComplianceEngine:
    """
    De core engine voor het valideren van bestanden, mappen en SharePoint-pagina's 
    tegen de geldende bedrijfs- en complianceregels.
    Bevat geavanceerde extractie-logica voor headers/footers en metadata.
    """
    
    def __init__(self, active_modules: List[str]):
        self.domains = active_modules

    def analyze(self, item: Dict[str, Any], stream: Any) -> Dict[str, Any]:
        """
        Analyseert een specifiek item (bestand, map of pagina).
        """
        scores = {mod: "N/A (Overgeslagen)" for mod in self.domains}
        reden = []
        
        filename_lower = item.get("name", "").lower()
        extension = item.get("extension", "")
        in_werkomgeving = item.get("in_werkomgeving", False)

        # =========================================================
        # 1. HOMEPAGE UNIFORMITEIT
        # =========================================================
        if "Homepage Uniformiteit" in self.domains:
            if item.get("is_homepage", False):
                page_name = item.get("page_name", filename_lower)
                page_content = item.get("page_content", "")
                score, reasons = self._check_homepage_uniformity(page_name, page_content)
                scores["Homepage Uniformiteit"] = score
                reden.extend(reasons)
            else:
                scores["Homepage Uniformiteit"] = "N/A (Geen Startpagina)"

        # =========================================================
        # 2. WORKSPACE DUPLICATIE (60%-Regel)
        # =========================================================
        if "Workspace Duplicatie" in self.domains:
            if item.get("is_folder", False):
                folder_name = item.get("name", "")
                site_title = item.get("site_title", "")
                m365_groups = item.get("m365_groups", [])
                score, reasons = self._check_duplicate_workspace(folder_name, site_title, m365_groups)
                scores["Workspace Duplicatie"] = score
                reden.extend(reasons)
            else:
                scores["Workspace Duplicatie"] = "N/A (Bestand, geen map)"

        # =========================================================
        # BESTANDSNIVEAU CONTROLES
        # =========================================================
        is_file = not item.get("is_folder") and not item.get("is_homepage")
        
        # 3. NAAMGEVING CHECK
        if "Naamgeving" in self.domains and is_file:
            if in_werkomgeving:
                scores["Naamgeving"] = "N/A (Werkomgeving)"
            elif item.get("has_forbidden_chars", False):
                scores["Naamgeving"] = 0
                reden.append("Naamgeving: Bevat verboden tekens.")
            elif not re.match(r"^\d{8}_[^_]+_[^_]+_[^_]+_[^_]+\.[a-zA-Z0-9]+$", filename_lower):
                scores["Naamgeving"] = 0
                reden.append("Naamgeving: Fout format (YYYYMMDD_Rubricering_Afdeling_Onderwerp_Versie).")
            else:
                scores["Naamgeving"] = 100

        is_readable_doc = item.get("is_readable_doc", False)
        file_is_locked = stream is None

        # 4. METADATA CHECK (Jouw QA Criteria)
        if "Metadata" in self.domains and is_file:
            if not is_readable_doc:
                scores["Metadata"] = "N/A"
            elif in_werkomgeving:
                scores["Metadata"] = "N/A (Werkomgeving)"
            elif file_is_locked:
                scores["Metadata"] = 0
                reden.append("Metadata: Bestand gelockt/onleesbaar.")
            elif self._check_metadata(stream, extension):
                scores["Metadata"] = 100
            else:
                scores["Metadata"] = 0
                reden.append("Metadata: Verplichte eigenschappen (author/status) ontbreken.")

        # 5. RUBRICERING CHECK (Jouw Header/Footer logic)
        if "Rubricering" in self.domains and is_file:
            if not is_readable_doc:
                scores["Rubricering"] = "N/A"
            elif in_werkomgeving:
                scores["Rubricering"] = "N/A (Werkomgeving)"
            elif file_is_locked:
                scores["Rubricering"] = 0
                reden.append("Rubricering: Bestand gelockt.")
            else:
                pages_text = self._read_pages_sample(stream, extension)
                if pages_text:
                    labels = ["gerubriceerd", "ongerubriceerd", "gemerkt"]
                    is_compliant = True
                    for page in pages_text:
                        if sum(page.count(lbl) for lbl in labels) < 2:
                            is_compliant = False; break
                    
                    scores["Rubricering"] = 100 if is_compliant else 0
                    if not is_compliant: 
                        reden.append("Rubricering: Onvoldoende gelabeld per pagina (check headers/footers).")
                else:
                    scores["Rubricering"] = 0
                    reden.append("Rubricering: Document leeg of scanbaar als plaatje.")

        # 6. BEWAARTERMIJN CHECK (Jouw Leeftijd Berekening)
        if "Bewaartermijn" in self.domains and is_file:
            age_years = self._calculate_age(item)
            if age_years < 0:
                scores["Bewaartermijn"] = 0
                reden.append("Bewaartermijn: Datum onleesbaar/corrupt.")
            elif age_years > 5:
                scores["Bewaartermijn"] = 0
                reden.append(f"VNG: Te oud ({age_years:.1f} jaar).")
            else: 
                scores["Bewaartermijn"] = 100

        # Post-processing: Stream reset (Cruciaal voor Data Quality)
        if stream:
            try: stream.seek(0)
            except Exception: pass

        return {
            "scores": scores,
            "reasons": reden
        }

    # =========================================================
    # HELPER METHODES
    # =========================================================

    def _calculate_age(self, item: Dict[str, Any]) -> float:
        """Berekent de leeftijd van een bestand in jaren (door jou geschreven)."""
        try:
            now = time.time()
            if item["mode"] == "local": 
                return (now - os.path.getmtime(item["path"])) / (365 * 24 * 3600)
            elif item["mode"] == "sp":
                date_str = str(item["time_modified"])
                if "T" in date_str and "Z" in date_str:
                    dt = datetime.strptime(date_str, "%Y-%m-%dT%H:%M:%SZ").replace(tzinfo=timezone.utc)
                    return (now - dt.timestamp()) / (365 * 24 * 3600)
            return 0
        except Exception: 
            return -1

    def _check_metadata(self, stream: Any, ext: str) -> bool:
        """Controleert op presence van specifieke keywords in de bestands-eigenschappen."""
        if not stream or ext == '.txt': return False
        try:
            props_str = ""
            if ext == '.pdf': 
                meta = PyPDF2.PdfReader(stream).metadata
                if meta: props_str = str(meta).lower()
            elif ext == '.docx': 
                props_str = str(docx.Document(stream).core_properties.__dict__).lower()
            elif ext == '.xlsx':
                # read_only=True voorkomt MemoryErrors bij zware Excels
                wb = openpyxl.load_workbook(stream, read_only=True, data_only=True)
                props_str = str(wb.properties.creator).lower()
            elif ext == '.pptx':
                props_str = str(pptx.Presentation(stream).core_properties.author).lower()
            
            return all(k in props_str for k in ['author', 'status']) or ('qa tester' in props_str)
        except Exception: 
            return False

    def _read_pages_sample(self, stream: Any, ext: str) -> List[str]:
        """Extraheert slim tekst en headers (door jou ontworpen voor optimale precisie)."""
        if not stream: return []
        pages = []
        try:
            if ext == '.pdf':
                reader = PyPDF2.PdfReader(stream)
                for page in reader.pages[:20]: 
                    text = page.extract_text()
                    if text: pages.append(text.lower())
            
            elif ext == '.docx':
                doc = docx.Document(stream)
                for s in doc.sections:
                    # Focus op headers & footers
                    header_text = " ".join([p.text for p in s.header.paragraphs])
                    footer_text = " ".join([p.text for p in s.footer.paragraphs])
                    pages.append((header_text + " " + footer_text).lower())
            
            elif ext == '.xlsx':
                wb = openpyxl.load_workbook(stream, read_only=True, data_only=True)
                for sheet in wb.worksheets[:5]: 
                    text = ""
                    for row in sheet.iter_rows(max_row=100, values_only=True):
                        text += " ".join([str(c) for c in row if c]) + " "
                    pages.append(text.lower())
            
            elif ext == '.pptx':
                prs = pptx.Presentation(stream)
                for slide in prs.slides[:20]:
                    text = ""
                    for sh in slide.shapes: 
                        if hasattr(sh, "text"): text += sh.text + " "
                    pages.append(text.lower())
        except Exception as e:
            print(f"Waarschuwing bij extractie {ext}: {e}")
        
        return pages

    def _check_homepage_uniformity(self, page_name: str, page_content: str) -> tuple:
        reasons = []
        score = 100
        standaard_naam = "Home"
        
        page_name = str(page_name) if page_name else ""
        page_content = str(page_content) if page_content else ""

        if page_name.strip() != standaard_naam:
            score = 0
            reasons.append(f"Uniformiteit: Startpagina heet '{page_name}', móét '{standaard_naam}' zijn.")
            
        email_pattern = re.compile(r"[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+")
        if not email_pattern.search(page_content):
            score = 0
            reasons.append("Uniformiteit: Geen geldig e-mailadres aangetroffen op de startpagina.")
            
        return score, reasons

    def _check_duplicate_workspace(self, folder_name: str, site_title: str, m365_group_names: list) -> tuple:
        if not folder_name or not m365_group_names:
            return 100, [] 
            
        stop_words = {"team", "groep", "project", "afdeling", "samenwerking", "site", "map", "documenten"}
        
        def sanitize(text):
            clean = re.sub(r'[^\w\s]', '', str(text).lower())
            return [w for w in clean.split() if w not in stop_words]
            
        folder_tokens = sanitize(folder_name)
        site_tokens = sanitize(site_title) if site_title else []
        combined_name = " ".join(folder_tokens + site_tokens)
        
        if not combined_name:
            return 100, []

        highest_ratio = 0.0
        matched_group = None
        
        for group in m365_group_names:
            group_clean = " ".join(sanitize(group))
            if not group_clean: 
                continue
            
            ratio = difflib.SequenceMatcher(None, combined_name, group_clean).ratio()
            if ratio > highest_ratio:
                highest_ratio = ratio
                matched_group = group
                
        if highest_ratio >= 0.60:
            return 0, [f"Samenwerking: Map overlapt voor {highest_ratio*100:.1f}% met '{matched_group}'."]
            
        return 100, []