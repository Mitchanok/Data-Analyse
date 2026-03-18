import os
import random
import time
from datetime import datetime, timedelta
import docx
import openpyxl
import pptx

# Configuratie
BASE_DIR = "Test_Compliance_Data"
SUB_DIRS = ["HR_Documenten", "Financien", "IT_Archief", "Tijdelijke_Downloads"]

EXTENSIONS_RISKY = [".exe", ".bat", ".ps1"]
EXTENSIONS_UNSUPPORTED = [".zip", ".mp4", ".iso"] 
FORBIDDEN_CHARS = [" ", "!", "@", "+"]

def setup_directories():
    if not os.path.exists(BASE_DIR):
        os.makedirs(BASE_DIR)
    for sub in SUB_DIRS:
        path = os.path.join(BASE_DIR, sub)
        if not os.path.exists(path):
            os.makedirs(path)

def spoof_file_age(filepath, years_old):
    """Past de wijzigingsdatum van een bestand aan om de VNG-bewaartermijn te testen."""
    if years_old > 0:
        past_time = time.time() - (years_old * 365 * 24 * 3600)
        os.utime(filepath, (past_time, past_time))

# ================= DOCUMENT GENERATOREN =================

def create_docx(filepath, add_metadata=True, add_rubricering=True):
    doc = docx.Document()
    
    if add_metadata:
        # Diepe XML Core Properties injectie
        doc.core_properties.author = "QA Tester"
        doc.core_properties.content_status = "Definitief"
        doc.core_properties.comments = "Bewaartermijn: VNG 5 Jaar"
        doc.core_properties.title = "Compliance Test Document"
        
    if add_rubricering:
        section = doc.sections[0]
        header = section.header
        header.paragraphs[0].text = "gerubriceerd gerubriceerd gerubriceerd"
        doc.add_paragraph("Dit document bevat zeer geheime VNG informatie.")
    else:
        doc.add_paragraph("Dit is een standaard document zonder geldige compliance labels.")
        
    doc.save(filepath)

def create_xlsx(filepath, add_metadata=True):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Financiele Data"
    
    if add_metadata:
        # Excel specifieke metadata
        wb.properties.creator = "QA Tester"
        wb.properties.title = "Compliance Test Excel"
        wb.properties.description = "Status: Definitief | Bewaartermijn: VNG 7 Jaar"
    
    # Voeg rubricering toe in de cellen
    ws['A1'] = "gerubriceerd"
    ws['A2'] = "gerubriceerd"
    ws['B1'] = "Geheime financiële cijfers 2023"
    
    wb.save(filepath)

def create_pptx(filepath, add_metadata=True):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(slide_layout)
    
    if add_metadata:
        # PowerPoint specifieke metadata
        prs.core_properties.author = "QA Tester"
        prs.core_properties.content_status = "Definitief"
        prs.core_properties.comments = "Bewaartermijn: VNG 10 Jaar"
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "gerubriceerd gerubriceerd"
    subtitle.text = "Geheime Presentatie"
    
    prs.save(filepath)

def create_fake_file(filepath, size_kb=1):
    with open(filepath, "wb") as f:
        f.write(os.urandom(size_kb * 1024))

def create_massive_sparse_file(filepath, size_gb):
    size_bytes = int(size_gb * 1024 * 1024 * 1024)
    with open(filepath, "wb") as f:
        f.seek(size_bytes - 1)
        f.write(b"\0")

# ================= TESTDATA GENERATIE =================

def generate_test_data():
    setup_directories()
    bestanden_gemaakt = 0
    print(f"🚀 Start generatie van geavanceerde testdata in map: ./{BASE_DIR} ...")

    # SCENARIO 1: Perfect Compliant Bestanden (Mix van Word, Excel, PPT) (~15 stuks)
    for i in range(5):
        date_str = (datetime.now() - timedelta(days=random.randint(1, 300))).strftime("%Y%m%d")
        
        # Maak een Word bestand
        create_docx(os.path.join(BASE_DIR, "HR_Documenten", f"{date_str}_Geheim_HR_Contract_v{i}.docx"), add_metadata=True, add_rubricering=True)
        # Maak een Excel bestand
        create_xlsx(os.path.join(BASE_DIR, "Financien", f"{date_str}_Geheim_Financien_Balans_v{i}.xlsx"), add_metadata=True)
        # Maak een PowerPoint bestand
        create_pptx(os.path.join(BASE_DIR, "IT_Archief", f"{date_str}_Geheim_IT_Presentatie_v{i}.pptx"), add_metadata=True)
        
        bestanden_gemaakt += 3

    # SCENARIO 2: Naamgeving Fails (~5 stuks)
    for i in range(5):
        bad_char = random.choice(FORBIDDEN_CHARS)
        filepath = os.path.join(BASE_DIR, "Financien", f"20230101_Geheim_Financien{bad_char}Rapportage v{i}.docx")
        create_docx(filepath, add_metadata=True, add_rubricering=True)
        bestanden_gemaakt += 1

    # SCENARIO 3: Inhoudelijke Fails (Bewust GEEN Metadata & Rubricering!) (~5 stuks)
    for i in range(5):
        date_str = datetime.now().strftime("%Y%m%d")
        filepath = os.path.join(BASE_DIR, "HR_Documenten", f"{date_str}_Ongerubriceerd_HR_Beleid_v{i}.docx")
        create_docx(filepath, add_metadata=False, add_rubricering=False)
        bestanden_gemaakt += 1

    # SCENARIO 4: Bewaartermijn Fails (Te oud > 5 jaar) (~5 stuks)
    for i in range(5):
        filepath = os.path.join(BASE_DIR, "HR_Documenten", f"20180512_Geheim_HR_OudDossier_v{i}.docx")
        create_docx(filepath, add_metadata=True, add_rubricering=True)
        spoof_file_age(filepath, years_old=random.uniform(5.5, 8.0)) 
        bestanden_gemaakt += 1

    # SCENARIO 5: Security & Foute Locaties (The 'Ugly') (~6 stuks)
    for i in range(3):
        ext = random.choice(EXTENSIONS_RISKY)
        create_fake_file(os.path.join(BASE_DIR, "Tijdelijke_Downloads", f"malware_test_{i}{ext}"), size_kb=5)
        bestanden_gemaakt += 1

    for i in range(3):
        ext = random.choice(EXTENSIONS_UNSUPPORTED)
        create_fake_file(os.path.join(BASE_DIR, "Tijdelijke_Downloads", f"vakantie_video_{i}{ext}"), size_kb=150)
        bestanden_gemaakt += 1

    # SCENARIO 6: Extreem grote bestanden (>2GB) 
    print("Maken van de 2GB+ testbestanden (Dit kost geen echte opslagruimte)...")
    create_massive_sparse_file(os.path.join(BASE_DIR, "Tijdelijke_Downloads", "enorme_bedrijfsfilm_2023.mp4"), 2.8)
    bestanden_gemaakt += 1
    
    date_str = datetime.now().strftime("%Y%m%d")
    create_massive_sparse_file(os.path.join(BASE_DIR, "IT_Archief", f"{date_str}_Geheim_IT_GroteDatadump_v1.zip"), 3.2)
    bestanden_gemaakt += 1

    print(f"✅ Succes! {bestanden_gemaakt} diverse testbestanden gegenereerd in de map '{BASE_DIR}'.")
    print("Test nu je ComplianceApp door de map 'Test_Compliance_Data' erin te slepen!")

if __name__ == "__main__":
    generate_test_data()