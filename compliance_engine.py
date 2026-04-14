import os
import re
import time
from datetime import datetime, timezone

import PyPDF2
import docx
import openpyxl
import pptx


class ComplianceEngine:
    """Compliance checks zonder overlap met CentraleEngine en KwaliteitEngine."""

    SUPPORTED_DOMAINS = {"Metadata", "Rubricering", "Bewaartermijn"}

    LABEL_GROUPS = {
        "ongerubriceerd": {"ongerubriceerd", "unclassified", "openbaar", "public"},
        "intern": {"intern", "internal", "internal use only"},
        "vertrouwelijk": {"vertrouwelijk", "confidential"},
        "geheim": {"geheim", "secret"},
        "gemerkt": {"gemerkt", "marked", "gemarkeerd"},
        "gerubriceerd": {"gerubriceerd", "classified"},
    }

    RETENTION_SIGNAL_WORDS = {
        "archief",
        "archived",
        "archive",
        "afgesloten",
        "gesloten",
        "final",
        "definitief",
        "signed",
        "getekend",
    }

    def __init__(self, active_modules):
        self.domains = [module for module in active_modules if module in self.SUPPORTED_DOMAINS]

    def analyze(self, item, stream):
        scores = {module: "N/A (Overgeslagen)" for module in self.domains}
        reasons = []

        extension = item["extension"]
        in_werkomgeving = item.get("in_werkomgeving", False)
        is_readable_doc = item.get("is_readable_doc", False)
        file_is_locked = stream is None

        metadata = self._extract_metadata(stream, extension) if stream and is_readable_doc else {}
        pages_text = None

        if "Metadata" in self.domains:
            score, module_reasons = self._check_metadata_rule(
                is_readable_doc=is_readable_doc,
                in_werkomgeving=in_werkomgeving,
                file_is_locked=file_is_locked,
                metadata=metadata,
            )
            scores["Metadata"] = score
            reasons.extend(module_reasons)

        if "Rubricering" in self.domains:
            if not file_is_locked and is_readable_doc:
                pages_text = self._read_pages_sample(stream, extension)
            score, module_reasons = self._check_rubricering_rule(
                is_readable_doc=is_readable_doc,
                in_werkomgeving=in_werkomgeving,
                file_is_locked=file_is_locked,
                pages_text=pages_text or [],
            )
            scores["Rubricering"] = score
            reasons.extend(module_reasons)

        if "Bewaartermijn" in self.domains:
            score, module_reasons = self._check_retention_rule(
                item=item,
                metadata=metadata,
                in_werkomgeving=in_werkomgeving,
            )
            scores["Bewaartermijn"] = score
            reasons.extend(module_reasons)

        if stream:
            try:
                stream.seek(0)
            except Exception:
                pass

        return {"scores": scores, "reasons": reasons}

    def _check_metadata_rule(self, is_readable_doc, in_werkomgeving, file_is_locked, metadata):
        if not is_readable_doc:
            return "N/A", []
        if in_werkomgeving:
            return "N/A (Werkomgeving)", []
        if file_is_locked:
            return 0, ["Metadata: bestand is gelockt of niet leesbaar."]

        owner_fields = {"author", "creator", "last_modified_by"}
        descriptive_fields = {"title", "subject", "keywords", "description", "comments"}
        lifecycle_fields = {"status", "category", "content_status"}

        owner_present = any(metadata.get(field) for field in owner_fields)
        descriptive_present = any(metadata.get(field) for field in descriptive_fields)
        lifecycle_present = any(metadata.get(field) for field in lifecycle_fields)

        if owner_present and (descriptive_present or lifecycle_present):
            return 100, []

        if owner_present or descriptive_present or lifecycle_present:
            missing = []
            if not owner_present:
                missing.append("eigenaar/auteur")
            if not descriptive_present and not lifecycle_present:
                missing.append("beschrijving of status")
            return 50, [f"Metadata: deels gevuld; ontbreekt {', '.join(missing)}."]

        return 0, ["Metadata: geen bruikbare auteur-, beschrijvende of statusmetadata gevonden."]

    def _check_rubricering_rule(self, is_readable_doc, in_werkomgeving, file_is_locked, pages_text):
        if not is_readable_doc:
            return "N/A", []
        if in_werkomgeving:
            return "N/A (Werkomgeving)", []
        if file_is_locked:
            return 0, ["Rubricering: bestand is gelockt of niet leesbaar."]
        if not pages_text:
            return 0, ["Rubricering: geen uitleesbare tekst gevonden voor controle."]

        detected_per_page = [self._detect_labels(text) for text in pages_text if text.strip()]
        if not detected_per_page:
            return 0, ["Rubricering: document bevat geen uitleesbare classificatietekst."]

        labeled_pages = [labels for labels in detected_per_page if labels]
        if not labeled_pages:
            return 0, ["Rubricering: geen herkenbare rubriceringslabels gevonden."]

        if len(labeled_pages) != len(detected_per_page):
            return 50, ["Rubricering: classificatie ontbreekt op een of meer pagina's of secties."]

        primary_labels = []
        for labels in labeled_pages:
            page_specific = sorted(label for label in labels if label not in {"gemerkt", "gerubriceerd"})
            if page_specific:
                primary_labels.append(page_specific[0])

        if primary_labels and len(set(primary_labels)) > 1:
            return 50, ["Rubricering: meerdere classificatieniveaus aangetroffen binnen hetzelfde document."]

        return 100, []

    def _check_retention_rule(self, item, metadata, in_werkomgeving):
        if in_werkomgeving:
            return "N/A (Werkomgeving)", []

        retention_until = self._parse_date(
            metadata.get("retention_until")
            or metadata.get("expiration_date")
            or metadata.get("expiry_date")
        )
        if retention_until:
            now = datetime.now(tz=retention_until.tzinfo or timezone.utc)
            days_remaining = (retention_until - now).days
            if days_remaining < 0:
                return 0, [f"Bewaartermijn: retentiedatum verstreken op {retention_until.date().isoformat()}."]
            if days_remaining <= 180:
                return 50, [f"Bewaartermijn: retentiedatum verloopt binnen {days_remaining} dagen."]
            return 100, []

        status_value = str(metadata.get("status", "")).lower()
        filename_value = str(item.get("name", "")).lower()
        has_retention_signal = status_value in self.RETENTION_SIGNAL_WORDS or any(
            word in filename_value for word in self.RETENTION_SIGNAL_WORDS
        )

        if not has_retention_signal:
            return "N/A (Geen retentie-indicatie)", []

        age_years = self._calculate_age(item)
        if age_years < 0:
            return 50, ["Bewaartermijn: leeftijd kon niet worden bepaald."]
        if age_years > 5:
            return 0, [f"Bewaartermijn: gearchiveerd of afgesloten document is {age_years:.1f} jaar oud."]
        if age_years > 4:
            return 50, [f"Bewaartermijn: document nadert de retentiegrens met {age_years:.1f} jaar leeftijd."]
        return 100, []

    def _calculate_age(self, item):
        try:
            now = time.time()
            if item["mode"] == "local":
                modified = item.get("time_modified", os.path.getmtime(item["path"]))
                return (now - modified) / (365 * 24 * 3600)
            if item["mode"] == "sp":
                modified_dt = self._parse_date(item.get("time_modified"))
                if modified_dt:
                    return (now - modified_dt.timestamp()) / (365 * 24 * 3600)
            return 0
        except Exception:
            return -1

    def _extract_metadata(self, stream, ext):
        if not stream:
            return {}

        try:
            stream.seek(0)
        except Exception:
            return {}

        metadata = {}
        try:
            if ext == ".pdf":
                pdf_meta = PyPDF2.PdfReader(stream).metadata or {}
                metadata = {
                    "author": self._clean_meta(pdf_meta.get("/Author")),
                    "creator": self._clean_meta(pdf_meta.get("/Creator")),
                    "title": self._clean_meta(pdf_meta.get("/Title")),
                    "subject": self._clean_meta(pdf_meta.get("/Subject")),
                    "keywords": self._clean_meta(pdf_meta.get("/Keywords")),
                    "status": self._clean_meta(pdf_meta.get("/Status")),
                    "category": self._clean_meta(pdf_meta.get("/Category")),
                    "content_status": self._clean_meta(pdf_meta.get("/ContentStatus")),
                    "description": self._clean_meta(pdf_meta.get("/Description")),
                    "retention_until": self._clean_meta(pdf_meta.get("/RetentionUntil")),
                    "expiry_date": self._clean_meta(pdf_meta.get("/ExpiryDate")),
                }
            elif ext == ".docx":
                props = docx.Document(stream).core_properties
                metadata = self._office_properties_to_dict(props)
            elif ext == ".xlsx":
                wb = openpyxl.load_workbook(stream, read_only=True)
                metadata = self._office_properties_to_dict(wb.properties)
                wb.close()
            elif ext == ".pptx":
                props = pptx.Presentation(stream).core_properties
                metadata = self._office_properties_to_dict(props)
        except Exception:
            metadata = {}
        finally:
            try:
                stream.seek(0)
            except Exception:
                pass

        return {key: value for key, value in metadata.items() if value}

    def _office_properties_to_dict(self, props):
        fields = {
            "author": getattr(props, "author", None),
            "creator": getattr(props, "creator", None),
            "last_modified_by": getattr(props, "last_modified_by", None),
            "title": getattr(props, "title", None),
            "subject": getattr(props, "subject", None),
            "keywords": getattr(props, "keywords", None),
            "description": getattr(props, "comments", None) or getattr(props, "description", None),
            "category": getattr(props, "category", None),
            "status": getattr(props, "status", None),
            "content_status": getattr(props, "content_status", None),
            "retention_until": getattr(props, "retention_until", None),
            "expiry_date": getattr(props, "expiry_date", None),
            "expiration_date": getattr(props, "expiration_date", None),
        }
        cleaned = {}
        for key, value in fields.items():
            cleaned_value = self._clean_meta(value)
            if cleaned_value:
                cleaned[key] = cleaned_value
        return cleaned

    def _clean_meta(self, value):
        if value is None:
            return None
        if isinstance(value, datetime):
            return value.isoformat()
        value = str(value).strip()
        if not value or value.lower() in {"none", "n/a", "unknown"}:
            return None
        return value

    def _read_pages_sample(self, stream, ext):
        if not stream:
            return []

        pages = []
        try:
            stream.seek(0)
            if ext == ".pdf":
                reader = PyPDF2.PdfReader(stream)
                for page in reader.pages[:20]:
                    text = page.extract_text() or ""
                    if text.strip():
                        pages.append(text.lower())
            elif ext == ".docx":
                doc = docx.Document(stream)
                header_parts = []
                footer_parts = []
                body_parts = []
                for section in doc.sections:
                    header_parts.extend(p.text for p in section.header.paragraphs if p.text.strip())
                    footer_parts.extend(p.text for p in section.footer.paragraphs if p.text.strip())
                body_parts.extend(p.text for p in doc.paragraphs[:200] if p.text.strip())
                combined = " ".join(header_parts + footer_parts + body_parts).lower().strip()
                if combined:
                    pages.append(combined)
            elif ext == ".xlsx":
                wb = openpyxl.load_workbook(stream, read_only=True)
                for sheet in wb.worksheets[:5]:
                    text_parts = [sheet.title]
                    for row in sheet.iter_rows(max_row=100, values_only=True):
                        text_parts.extend(str(cell) for cell in row if cell not in (None, ""))
                    sheet_text = " ".join(text_parts).lower().strip()
                    if sheet_text:
                        pages.append(sheet_text)
                wb.close()
            elif ext == ".pptx":
                prs = pptx.Presentation(stream)
                for slide in prs.slides[:20]:
                    parts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and str(shape.text).strip():
                            parts.append(shape.text)
                    slide_text = " ".join(parts).lower().strip()
                    if slide_text:
                        pages.append(slide_text)
            elif ext == ".txt":
                text = stream.read().decode("utf-8", errors="ignore").lower()
                chunks = [text[i:i + 4000].strip() for i in range(0, len(text), 4000)]
                pages = [chunk for chunk in chunks[:20] if chunk]
        except Exception:
            pages = []
        finally:
            try:
                stream.seek(0)
            except Exception:
                pass

        return pages

    def _detect_labels(self, text):
        found = set()
        normalized = re.sub(r"\s+", " ", text.lower())
        for canonical_label, variants in self.LABEL_GROUPS.items():
            if any(variant in normalized for variant in variants):
                found.add(canonical_label)
        return found

    def _parse_date(self, value):
        if not value:
            return None
        if isinstance(value, datetime):
            return value if value.tzinfo else value.replace(tzinfo=timezone.utc)
        if hasattr(value, "timestamp"):
            try:
                return datetime.fromtimestamp(value.timestamp(), tz=timezone.utc)
            except Exception:
                pass

        value_str = str(value).strip()
        formats = (
            "%Y-%m-%d",
            "%Y-%m-%dT%H:%M:%S",
            "%Y-%m-%dT%H:%M:%SZ",
            "%Y-%m-%d %H:%M:%S",
            "%d-%m-%Y",
            "%d/%m/%Y",
        )
        for fmt in formats:
            try:
                parsed = datetime.strptime(value_str, fmt)
                if parsed.tzinfo is None:
                    parsed = parsed.replace(tzinfo=timezone.utc)
                return parsed
            except ValueError:
                continue
        return None
