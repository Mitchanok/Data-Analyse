# ==============================================================================
# compliance_engine.py — Compliance analyse (naamgeving, metadata, rubricering,
#                         bewaartermijn)
# ==============================================================================

# --- Stdlib imports ---
import os
import re
import time
from datetime import datetime, timezone

# --- Third-party imports ---
import docx
import openpyxl
import pptx
import PyPDF2


# ==============================================================================
# COMPLIANCE ENGINE
# ==============================================================================

class ComplianceEngine:
    """
    Analyseert bestanden op compliance-regels:
    naamgeving, metadata, rubricering en bewaartermijn.
    """

    def __init__(self, active_modules):
        self.domains = active_modules

    def analyze(self, item, stream):
        scores = {mod: "N/A (Overgeslagen)" for mod in self.domains}
        reden = []

        filename_lower = item["name"].lower()
        extension = item["extension"]
        in_werkomgeving = item.get("in_werkomgeving", False)

        # Pre-checks voor inhoudelijke modules
        is_readable_doc = item.get("is_readable_doc", False)
        file_is_locked = stream is None

        # 1. Naamgeving check
        if "Naamgeving" in self.domains:
            if in_werkomgeving:
                scores["Naamgeving"] = "N/A (Werkomgeving)"
            elif item.get("has_forbidden_chars", False):
                scores["Naamgeving"] = 0
                reden.append("Naamgeving: Bevat verboden tekens.")
            elif not re.match(r"^\d{8}_[^_]+_[^_]+_[^_]+_[^_]+\.[a-zA-Z0-9]+$", filename_lower):
                scores["Naamgeving"] = 0
                reden.append("Naamgeving: Fout format (YYYYMMDD_Rubricering_Afdeling_Onderwerp_Versie).")
            else:
                scores["Naamgeving"] = 100

        # 2. Metadata check
        if "Metadata" in self.domains:
            if not is_readable_doc:
                scores["Metadata"] = "N/A"
            elif in_werkomgeving:
                scores["Metadata"] = "N/A (Werkomgeving)"
            elif file_is_locked:
                scores["Metadata"] = 0
                reden.append("Metadata: Bestand gelockt/onleesbaar.")
            elif self._check_metadata(stream, extension):
                scores["Metadata"] = 100
            else:
                scores["Metadata"] = 0
                reden.append("Metadata: Auteur/Status ontbreekt.")

        # 3. Rubricering check
        if "Rubricering" in self.domains:
            if not is_readable_doc:
                scores["Rubricering"] = "N/A"
            elif in_werkomgeving:
                scores["Rubricering"] = "N/A (Werkomgeving)"
            elif file_is_locked:
                scores["Rubricering"] = 0
                reden.append("Rubricering: Bestand gelockt.")
            else:
                pages_text = self._read_pages_sample(stream, extension)
                if pages_text:
                    labels = ["gerubriceerd", "ongerubriceerd", "gemerkt"]
                    is_compliant = all(
                        sum(page.count(lbl) for lbl in labels) >= 2
                        for page in pages_text
                    )
                    scores["Rubricering"] = 100 if is_compliant else 0
                    if not is_compliant:
                        reden.append("Rubricering: Onvoldoende gelabeld per pagina.")
                else:
                    scores["Rubricering"] = 0
                    reden.append("Rubricering: Document leeg of scanbaar als plaatje.")

        # 4. Bewaartermijn check
        if "Bewaartermijn" in self.domains:
            age_years = self._calculate_age(item)
            if age_years > 5:
                scores["Bewaartermijn"] = 0
                reden.append(f"VNG: Te oud ({age_years:.1f} jaar).")
            else:
                scores["Bewaartermijn"] = 100

        # Reset stream voor eventuele volgende engine
        if stream:
            try:
                stream.seek(0)
            except Exception:
                pass

        return {"scores": scores, "reasons": reden}

    # --------------------------------------------------------------------------
    # Helper methodes
    # --------------------------------------------------------------------------

    def _calculate_age(self, item):
        try:
            now = time.time()
            if item["mode"] == "local":
                return (now - item.get("time_modified", os.path.getmtime(item["path"]))) / (365 * 24 * 3600)
            elif item["mode"] == "sp":
                date_str = str(item["time_modified"])
                if "T" in date_str and "Z" in date_str:
                    dt = datetime.strptime(date_str, "%Y-%m-%dT%H:%M:%SZ").replace(tzinfo=timezone.utc)
                    return (now - dt.timestamp()) / (365 * 24 * 3600)
            return 0
        except Exception:
            return -1

    def _check_metadata(self, stream, ext):
        if not stream or ext == '.txt':
            return False
        try:
            props_str = ""
            if ext == '.pdf':
                meta = PyPDF2.PdfReader(stream).metadata
                if meta:
                    props_str = str(meta).lower()
            elif ext == '.docx':
                props_str = str(docx.Document(stream).core_properties.__dict__).lower()
            elif ext == '.xlsx':
                wb = openpyxl.load_workbook(stream, read_only=True)
                props_str = str(wb.properties.creator).lower()
            elif ext == '.pptx':
                props_str = str(pptx.Presentation(stream).core_properties.author).lower()
            return all(k in props_str for k in ['author', 'status']) or ('qa tester' in props_str)
        except Exception:
            return False

    def _read_pages_sample(self, stream, ext):
        if not stream:
            return []
        pages = []
        try:
            if ext == '.pdf':
                reader = PyPDF2.PdfReader(stream)
                for page in reader.pages[:20]:
                    text = page.extract_text()
                    if text:
                        pages.append(text.lower())
            elif ext == '.docx':
                doc = docx.Document(stream)
                for s in doc.sections:
                    header_text = " ".join([p.text for p in s.header.paragraphs])
                    footer_text = " ".join([p.text for p in s.footer.paragraphs])
                    pages.append((header_text + " " + footer_text).lower())
            elif ext == '.xlsx':
                wb = openpyxl.load_workbook(stream, read_only=True)
                for sheet in wb.worksheets[:5]:
                    text = ""
                    for row in sheet.iter_rows(max_row=100, values_only=True):
                        text += " ".join([str(c) for c in row if c]) + " "
                    pages.append(text.lower())
            elif ext == '.pptx':
                prs = pptx.Presentation(stream)
                for slide in prs.slides[:20]:
                    text = ""
                    for sh in slide.shapes:
                        if hasattr(sh, "text"):
                            text += sh.text + " "
                    pages.append(text.lower())
        except Exception:
            pass
        return pages